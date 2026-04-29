from pptx import Presentation

prs = Presentation('Batch-17_ppt.pptx')

print(f"{'='*80}")
print(f"PRESENTATION ANALYSIS - {len(prs.slides)} Slides")
print(f"{'='*80}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*80}")
    print(f"SLIDE {i}")
    print(f"{'='*80}")
    
    # Extract all text from shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            # Print with line breaks preserved
            print(text)
            print("-" * 40)
    
    # Count images
    image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)  # 13 = Picture
    if image_count > 0:
        print(f"[{image_count} image(s) on this slide]")

print(f"\n{'='*80}")
print("ANALYSIS COMPLETE")
print(f"{'='*80}")
